from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Define brand colors
PRIMARY_PURPLE = RGBColor(77, 58, 168)
ACCENT_PURPLE = RGBColor(117, 87, 255)
ICE_BLUE = RGBColor(0, 212, 255)
PENGUIN_ORANGE = RGBColor(255, 107, 53)
GNOME_RED = RGBColor(255, 51, 102)
DARK_BG = RGBColor(26, 26, 46)
LIGHT_TEXT = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Style title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_PURPLE
    
    # Style subtitle
    subtitle_frame = subtitle_shape.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = DARK_BG
    
    return slide

def add_content_slide(prs, title, content_items, level=0):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_PURPLE
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = level
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_BG
    
    return slide

def add_section_slide(prs, title, subtitle=""):
    """Add a section divider slide"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title box
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_PURPLE
    
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.5), width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.color.rgb = ACCENT_PURPLE
    
    # Add colored bar at bottom
    left = Inches(0)
    top = Inches(6.5)
    width = Inches(10)
    height = Inches(0.2)
    
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICE_BLUE
    shape.line.fill.background()
    
    return slide

# Slide 1: Title
add_title_slide(prs, "THE LIGHT PARK 2026", "Creative Playbook & Strategic Roadmap")

# Slide 2: Table of Contents
add_content_slide(prs, "Table of Contents", [
    "Executive Summary",
    "Project Timeline Overview",
    "February 2026 - Foundation & Character Development",
    "March 2026 - Content Creation & Digital Platforms",
    "April 2026 - Integration & Finalization",
    "May 2026 - Public Launch & Production",
    "June 2026 - Community Engagement",
    "Concept Visualizations"
])

# Slide 3: Executive Summary
add_section_slide(prs, "Executive Summary", "Strategic Vision for 2026 Season")

# Slide 4: Mission Statement
add_content_slide(prs, "Mission Statement", [
    "THE LIGHT PARK 2026 season represents a transformative leap forward in holiday entertainment",
    "Expanding to 8 venues across multiple locations",
    "Introducing DJ Polar Ice and Pixel Penguin character-driven experiences",
    "Integrating cutting-edge technology and immersive storytelling",
    "Creating unforgettable holiday memories for all guests"
])

# Slide 5: Key Objectives
add_content_slide(prs, "Key Objectives", [
    "✓ Venue Expansion: Launch new Oklahoma City park location",
    "✓ Character Integration: Fully integrate DJ Polar Ice and Pixel Penguin",
    "✓ Experience Enhancement: Redesign tent experiences with digital screens",
    "✓ Digital Transformation: Launch mobile app and YouTube channel",
    "✓ Content Creation: Develop new audio sequences and theme songs"
])

# Slide 6: Brand Cohesion
add_content_slide(prs, "Brand Cohesion", [
    "Strict adherence to THE LIGHT PARK brand identity",
    "Signature color palette: Primary Purple, Ice Blue, Accent Purple, Gnome Red",
    "Character universe consistency across all touchpoints",
    "Professional, polished presentation for stakeholder review",
    "Cohesive brand presence across all materials and experiences"
])

# Slide 7: Timeline Section
add_section_slide(prs, "Project Timeline Overview", "February - June 2026")

# Slide 8: Gantt Chart Overview
add_content_slide(prs, "Phase Overview - Gantt Chart", [
    "February: Character Development & Venue Design",
    "  • DJ Polar Ice avatar development",
    "  • Pixel Penguin concept design",
    "  • Oklahoma City park layout",
    "  • Tent enhancement concepts",
    "",
    "March: Content Creation & Digital Launch",
    "  • YouTube channel setup",
    "  • Festival map prototypes",
    "  • Mobile app development",
    "  • Audio track production"
])

# Slide 9: Linear Timeline
add_content_slide(prs, "Key Milestones - Linear Timeline", [
    "February 2026: Foundation & Character Development",
    "  ✓ DJ Polar Ice avatar photo shoot complete",
    "  ✓ Pixel Penguin designs finalized",
    "  ✓ OKC park layout approved",
    "",
    "March 2026: Content & Digital Platform Launch",
    "  ✓ YouTube channel created (private)",
    "  ✓ Festival map prototypes developed",
    "  ✓ Mobile app beta testing begins",
    "",
    "April 2026: Integration & Finalization",
    "  ✓ Bumper tracks integrated",
    "  ✓ Festival maps finalized",
    "  ✓ Branding approved"
])

# Slide 10: February Section
add_section_slide(prs, "February 2026", "Foundation Building & Character Development")

# Slide 11: DJ Polar Ice Development
add_content_slide(prs, "DJ Polar Ice Avatar Development", [
    "HIGH PRIORITY - Mascot Photo Shoot & Facial Animation",
    "  • Schedule professional photographer with mascot experience",
    "  • Capture multiple angles and expressions",
    "  • Create animation library for marketing and shows",
    "",
    "HIGH PRIORITY - Lip-Sync Modeling & Gesture Animation",
    "  • Create phoneme library (English/Spanish)",
    "  • Design gesture animations",
    "  • Test animation smoothness",
    "",
    "HIGH PRIORITY - Voice Creation",
    "  • Audition bilingual voice actors",
    "  • Record core phrases",
    "  • Integrate with show control systems"
])

# Slide 12: Pixel Penguin & Venue Design
add_content_slide(prs, "Pixel Penguin & Venue Design", [
    "HIGH PRIORITY - Pixel Penguin Character Design",
    "  • Develop 5-7 concept variations",
    "  • Focus on mischievous, tech-savvy traits",
    "  • Create final 3D model",
    "",
    "HIGH PRIORITY - Oklahoma City Park Layout",
    "  • Conduct site survey and measurement",
    "  • Design traffic flow patterns",
    "  • Map emergency exits and safety compliance",
    "  • Position concession tent and merchandise areas",
    "",
    "MEDIUM PRIORITY - Themed Area Development",
    "  • Audit existing prop inventory",
    "  • Design layouts for returning venues"
])

# Slide 13: Tent Enhancement
add_content_slide(prs, "Tent Experience Enhancement", [
    "HIGH PRIORITY - Tent Redesign",
    "  • Design 16' tall pixel screens for entrances",
    "  • Create custom rolling shelving units",
    "  • Design rolling cart for popcorn machines",
    "  • Develop standardized tent layout template",
    "",
    "HIGH PRIORITY - Merchandising Strategy",
    "  • Display DJ Polar Ice plushies",
    "  • Electric lemonade, cake pops, cotton candy",
    "  • Generate AI concept images",
    "",
    "MEDIUM PRIORITY - Marketing Assets Organization",
    "  • Create centralized digital asset library",
    "  • Establish file naming conventions"
])

# Slide 14: March Section
add_section_slide(prs, "March 2026", "Content Creation & Digital Platform Development")

# Slide 15: Digital Platforms
add_content_slide(prs, "Digital Platforms Launch", [
    "MEDIUM PRIORITY - YouTube Channel",
    "  • Create channel with proper branding",
    "  • Design channel art and thumbnails",
    "  • Upload initial content for review",
    "  • Conduct brand consistency audit",
    "",
    "HIGH PRIORITY - Festival Map Prototypes",
    "  • Design site-specific maps for all 8 venues",
    "  • Include DJ Polar Ice story teaser",
    "  • Create sponsorship placement zones",
    "  • Begin sponsor outreach"
])

# Slide 16: Mobile App & Content
add_content_slide(prs, "Mobile App & Content Creation", [
    "HIGH PRIORITY - Mobile App Development",
    "  • Provide branding materials to developer",
    "  • Review UI/UX designs",
    "  • Test beta versions",
    "  • Ensure brand consistency",
    "",
    "MEDIUM PRIORITY - Storybook Concept",
    "  • Write DJ Polar Ice origin story",
    "  • Create teaser excerpts",
    "  • Develop character dialogue",
    "",
    "HIGH PRIORITY - Theme Song/Jingle",
    "  • Compose signature audio branding",
    "  • Create instrumental and vocal versions",
    "  • Integrate into marketing and shows"
])

# Slide 17: Audio & Sequences
add_content_slide(prs, "Audio & Sequence Production", [
    "MEDIUM PRIORITY - Character Dialog",
    "  • Write bumper track scripts",
    "  • Develop bilingual content",
    "  • Add personality to show experience",
    "",
    "HIGH PRIORITY - New Audio Track Production",
    "  • Create minimum 3 new sequences",
    "  • Select current music tracks",
    "  • Design light synchronization patterns",
    "",
    "MEDIUM PRIORITY - Digital Character Images",
    "  • Create pixel screen-optimized graphics",
    "  • Design phoneme animations",
    "  • Test display quality",
    "",
    "HIGH PRIORITY - xLights Layout Completion",
    "  • Map sequences to all 8 venues",
    "  • Test HTML upload scripts",
    "  • Enable mid-season updates"
])

# Slide 18: April Section
add_section_slide(prs, "April 2026", "Integration & Finalization Phase")

# Slide 19: Integration Tasks
add_content_slide(prs, "Integration & Finalization", [
    "HIGH PRIORITY - Bumper Track Integration",
    "  • Import all tracks into xLights",
    "  • Add theme song backgrounds",
    "  • Maintain light animation during transitions",
    "",
    "HIGH PRIORITY - Site Layout Finalization",
    "  • Complete all 8 festival maps",
    "  • Obtain domestic printing quotes",
    "  • Solicit overseas printing quotes",
    "  • Make final printing decision",
    "",
    "HIGH PRIORITY - Branding Finalization",
    "  • Comprehensive graphic design review",
    "  • Ensure brand cohesion",
    "  • Obtain stakeholder approvals"
])

# Slide 20: App Finalization
add_content_slide(prs, "App Finalization", [
    "HIGH PRIORITY - Mobile App Completion",
    "  • Confirm all use cases",
    "  • Test ticketing system integration",
    "  • Conduct user acceptance testing",
    "  • Prepare for app store submission",
    "",
    "HIGH PRIORITY - Quality Assurance",
    "  • Test all features across devices",
    "  • Verify brand alignment",
    "  • Validate user experience",
    "  • Final bug fixes and improvements"
])

# Slide 21: May Section
add_section_slide(prs, "May 2026", "Public Launch & Production Phase")

# Slide 22: Social Media Launch
add_content_slide(prs, "Social Media Launch", [
    "HIGH PRIORITY - Multi-Platform Launch",
    "  • Begin content release on YouTube, TikTok, Instagram",
    "  • Post teaser content and character reveals",
    "  • Monitor engagement metrics daily",
    "  • Adjust strategy based on performance",
    "",
    "Key Success Metrics",
    "  • Viewer engagement and watch time",
    "  • Social media follower growth",
    "  • Content sharing and virality",
    "  • Brand awareness and recognition"
])

# Slide 23: Production & Sales
add_content_slide(prs, "Production & Early Bird Sales", [
    "HIGH PRIORITY - Tent Improvement Production",
    "  • Finalize concession enhancement designs",
    "  • Begin manufacturing pixel screens",
    "  • Procurement of rolling shelving units",
    "  • Fabrication of rolling carts",
    "",
    "MEDIUM PRIORITY - Early Bird Sales Planning",
    "  • Define promotional packages",
    "  • Establish pricing tiers",
    "  • Create marketing materials",
    "  • Set launch dates",
    "",
    "HIGH PRIORITY - Sequence Upload Process",
    "  • Upload new sequences to all props",
    "  • Update existing content",
    "  • Test playback quality",
    "  • Continue through July"
])

# Slide 24: June Section
add_section_slide(prs, "June 2026", "Community Engagement & Pre-Season Preparation")

# Slide 25: Community Outreach
add_content_slide(prs, "Community Engagement", [
    "MEDIUM PRIORITY - Community Outreach",
    "  • Contact local schools for spirit nights",
    "  • Develop school partnership packages",
    "  • Approach businesses for group sales",
    "  • Create promotional materials",
    "",
    "MEDIUM PRIORITY - App Audio Content",
    "  • Develop pre-arrival audio experience",
    "  • Create queue entertainment content",
    "  • Explore sponsorship opportunities",
    "  • Integrate into mobile app",
    "",
    "Partnership Goals",
    "  • Build community relationships",
    "  • Increase group ticket sales",
    "  • Enhance brand local presence"
])

# Slide 26: Concept Visualizations
add_section_slide(prs, "Concept Visualizations", "AI-Generated Images for Tent Enhancements")

# Slide 27: Concept Images
add_content_slide(prs, "Tent Enhancement Concepts", [
    "🎪 Grand Entrance Design",
    "  • 16' tall pixel screens at tent leg entrances",
    "  • DJ Polar Ice content displays",
    "  • Immersive arrival experience",
    "",
    "🛍️ Merchandise Display",
    "  • Custom rolling shelving units",
    "  • DJ Polar Ice plushies showcase",
    "  • Electric lemonade, cake pops, cotton candy",
    "  • Organized, visually appealing arrangement",
    "",
    "🍿 Concession Experience",
    "  • Rolling cart for popcorn machines",
    "  • Themed branding design",
    "  • Efficient service flow",
    "  • Standardized across venues"
])

# Slide 28: Closing
add_title_slide(prs, "Ready for 2026!", "THE LIGHT PARK - Creating Magical Holiday Experiences")

# Save presentation
prs.save('THE_LIGHT_PARK_2026_Playbook.pptx')
print("PowerPoint presentation created successfully!")